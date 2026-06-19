"""Загрузка реальной базы знаний из папки в единый контекст RLM.

Папка целиком склеивается в одну строку `context` (родная для RLM схема:
длинный контекст хранится как одна Python-переменная, корневая модель грепает/
режет его сама). Каждый файл предваряется разделителем с относительным путём,
структура (таблицы, листы, слайды, страницы) по возможности сохраняется —
это помогает grep/aggregation-стратегиям движка.

Поддерживаются: .docx, .xlsx, .csv, .pptx, .pdf, .txt, .md. Битые файлы и
отсутствующие библиотеки не роняют прогон — вместо содержимого пишется пометка.
"""

from __future__ import annotations

import csv as _csv
import os
from typing import Callable

from eval.datasets import Task

# Расширения, которые умеем читать (в нижнем регистре, с точкой).
SUPPORTED_EXTS = {".docx", ".xlsx", ".csv", ".pptx", ".pdf", ".txt", ".md"}


def _read_docx(path: str) -> str:
    from docx import Document  # python-docx

    doc = Document(path)
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for ti, table in enumerate(doc.tables):
        parts.append(f"--- Таблица {ti + 1} ---")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"--- Лист: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _read_csv(path: str) -> str:
    parts: list[str] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for row in _csv.reader(f):
            parts.append(" | ".join(row))
    return "\n".join(parts)


def _read_pptx(path: str) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(path)
    parts: list[str] = []
    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Слайд {si} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts: list[str] = []
    for pi, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"--- Стр. {pi} ---")
            parts.append(text)
    return "\n".join(parts)


def _read_text(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


_READERS: dict[str, Callable[[str], str]] = {
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".csv": _read_csv,
    ".pptx": _read_pptx,
    ".pdf": _read_pdf,
    ".txt": _read_text,
    ".md": _read_text,
}


def _iter_files(folder: str, recursive: bool) -> list[str]:
    """Отсортированный список поддерживаемых файлов папки (стабильный порядок)."""
    found: list[str] = []
    if recursive:
        for root, _dirs, files in os.walk(folder):
            for name in files:
                if os.path.splitext(name)[1].lower() in SUPPORTED_EXTS:
                    found.append(os.path.join(root, name))
    else:
        for name in os.listdir(folder):
            full = os.path.join(folder, name)
            if os.path.isfile(full) and os.path.splitext(name)[1].lower() in SUPPORTED_EXTS:
                found.append(full)
    return sorted(found)


def load_folder(folder: str, *, recursive: bool = True) -> tuple[str, int]:
    """Склеить все поддерживаемые файлы папки в один контекст.

    Возвращает (combined_text, n_files). Битые файлы/нет библиотеки → пометка
    в тексте, прогон не падает. n_files считает только реально прочитанные файлы.
    """
    folder = os.path.expanduser(folder)
    if not os.path.isdir(folder):
        raise FileNotFoundError(f"Папка не найдена: {folder}")

    files = _iter_files(folder, recursive)
    if not files:
        raise ValueError(
            f"В папке нет поддерживаемых файлов ({', '.join(sorted(SUPPORTED_EXTS))}): {folder}"
        )

    blocks: list[str] = []
    n_ok = 0
    for path in files:
        rel = os.path.relpath(path, folder)
        ext = os.path.splitext(path)[1].lower()
        header = f"\n\n===== ФАЙЛ: {rel} =====\n"
        try:
            body = _READERS[ext](path)
            n_ok += 1
        except ImportError as e:
            body = f"[пропущено: не установлена библиотека для {ext} — {e}]"
        except Exception as e:  # noqa: BLE001 — один битый файл не должен ронять прогон
            body = f"[не удалось прочитать файл ({ext}): {e}]"
        blocks.append(header + body)

    return "".join(blocks).strip(), n_ok


def build_real_tasks(rows: list[tuple[str, str]], *, recursive: bool = True) -> list[Task]:
    """Собрать задачи из строк (папка, вопрос) интерфейса.

    Пустые строки (нет ни папки, ни вопроса) пропускаются. Для каждой заполненной
    грузит папку и создаёт Task без эталона (answer_kind="none"). Если папка или
    вопрос заполнены лишь частично — поднимает ValueError с понятным текстом.
    """
    tasks: list[Task] = []
    for i, (folder, question) in enumerate(rows, 1):
        folder = (folder or "").strip()
        question = (question or "").strip()
        if not folder and not question:
            continue
        if not folder or not question:
            raise ValueError(
                f"Строка {i}: заполните и папку, и вопрос (или оставьте оба пустыми)."
            )
        context, n_files = load_folder(folder, recursive=recursive)
        label = os.path.basename(os.path.normpath(folder)) or f"kb{i}"
        tasks.append(Task(
            id=f"{label}_{i}",
            type="real",
            question=question,
            answer="",
            context=context,
            char_len=len(context),
            answer_kind="none",
            source=folder,
            n_files=n_files,
        ))
    return tasks
